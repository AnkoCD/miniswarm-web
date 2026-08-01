import csv
import copy
import os
import json
import re
import shutil
import subprocess
import sys
import tempfile
import threading
import time
import zipfile
from html.parser import HTMLParser
from pathlib import Path, PurePath
from typing import Any
from urllib.parse import urlparse

from runner_app.config import RunnerSettings, get_settings
from runner_app.agent_scope import RunnerScopeError, enforce_request_scope
from runner_app.schemas import ToolRequest, ToolResponse


class ToolRejected(ValueError):
    pass


PPT_SKILL_NAME = "guizang-ppt-skill"
URL_PATTERN = re.compile(r"https?://[^\s<>\]\[()\"']+", re.IGNORECASE)
_ANYSEARCH_EXTRACT_CACHE_TTL_SECONDS = 30 * 60
_ANYSEARCH_EXTRACT_CACHE_MAX_ITEMS = 128
_ANYSEARCH_EXTRACT_CACHE: dict[str, tuple[float, dict[str, Any]]] = {}
_ANYSEARCH_EXTRACT_LOCKS: dict[str, threading.Lock] = {}
_ANYSEARCH_EXTRACT_META_LOCK = threading.Lock()


def _anysearch_extract_cache_get(key: str) -> dict[str, Any] | None:
    now = time.monotonic()
    with _ANYSEARCH_EXTRACT_META_LOCK:
        entry = _ANYSEARCH_EXTRACT_CACHE.get(key)
        if entry is None:
            return None
        created_at, data = entry
        if now - created_at > _ANYSEARCH_EXTRACT_CACHE_TTL_SECONDS:
            _ANYSEARCH_EXTRACT_CACHE.pop(key, None)
            return None
        return copy.deepcopy(data)


def _anysearch_extract_cache_put(key: str, data: dict[str, Any]) -> None:
    with _ANYSEARCH_EXTRACT_META_LOCK:
        if len(_ANYSEARCH_EXTRACT_CACHE) >= _ANYSEARCH_EXTRACT_CACHE_MAX_ITEMS:
            oldest_key = min(
                _ANYSEARCH_EXTRACT_CACHE,
                key=lambda item: _ANYSEARCH_EXTRACT_CACHE[item][0],
            )
            _ANYSEARCH_EXTRACT_CACHE.pop(oldest_key, None)
        _ANYSEARCH_EXTRACT_CACHE[key] = (time.monotonic(), copy.deepcopy(data))


def _anysearch_extract_lock(key: str) -> threading.Lock:
    with _ANYSEARCH_EXTRACT_META_LOCK:
        return _ANYSEARCH_EXTRACT_LOCKS.setdefault(key, threading.Lock())


class _HTMLSummaryParser(HTMLParser):
    def __init__(self) -> None:
        super().__init__(convert_charrefs=True)
        self.title_parts: list[str] = []
        self.in_title = False
        self.text_chars = 0
        self.elements = 0
        self.scripts = 0
        self.forms = 0
        self.slides = 0

    def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        self.elements += 1
        normalized = tag.lower()
        if normalized == "title":
            self.in_title = True
        elif normalized == "script":
            self.scripts += 1
        elif normalized == "form":
            self.forms += 1
        classes = next((value or "" for key, value in attrs if key.lower() == "class"), "")
        if "slide" in classes.split():
            self.slides += 1

    def handle_endtag(self, tag: str) -> None:
        if tag.lower() == "title":
            self.in_title = False

    def handle_data(self, data: str) -> None:
        value = data.strip()
        if value:
            self.text_chars += len(value)
            if self.in_title:
                self.title_parts.append(value)


def _pptx_layout_issues(presentation) -> list[dict]:
    """Detect clear text overflow/overlap without flagging card backgrounds."""
    points = 12_700
    tolerance = 2 * points
    issues: list[dict] = []
    for slide_number, slide in enumerate(presentation.slides, 1):
        text_shapes: list[tuple[int, object, str]] = []
        for shape_index, shape in enumerate(slide.shapes):
            text = str(getattr(shape, "text", "") or "").strip()
            if not text or not all(
                isinstance(getattr(shape, name, None), int)
                for name in ("left", "top", "width", "height")
            ):
                continue
            text_shapes.append((shape_index, shape, text))
            if (
                shape.left < -tolerance
                or shape.top < -tolerance
                or shape.left + shape.width > presentation.slide_width + tolerance
                or shape.top + shape.height > presentation.slide_height + tolerance
            ):
                issues.append(
                    {
                        "slide": slide_number,
                        "kind": "text_out_of_bounds",
                        "shape": shape_index,
                        "text": text.replace("\n", " / ")[:120],
                    }
                )
        for first_index in range(len(text_shapes)):
            first_shape_index, first, first_text = text_shapes[first_index]
            for second_shape_index, second, second_text in text_shapes[first_index + 1:]:
                overlap_width = min(first.left + first.width, second.left + second.width) - max(first.left, second.left)
                overlap_height = min(first.top + first.height, second.top + second.height) - max(first.top, second.top)
                if overlap_width <= tolerance or overlap_height <= tolerance:
                    continue
                overlap_area = overlap_width * overlap_height
                smaller_area = min(first.width * first.height, second.width * second.height)
                if smaller_area <= 0 or overlap_area / smaller_area < 0.20:
                    continue
                issues.append(
                    {
                        "slide": slide_number,
                        "kind": "text_overlap",
                        "shapes": [first_shape_index, second_shape_index],
                        "texts": [
                            first_text.replace("\n", " / ")[:100],
                            second_text.replace("\n", " / ")[:100],
                        ],
                        "overlap_ratio": round(overlap_area / smaller_area, 3),
                    }
                )
        if len(issues) >= 20:
            break
    return issues[:20]


def _run_quality_command(
    command: list[str],
    *,
    cwd: Path,
    settings: RunnerSettings,
    timeout: int = 180,
    env: dict[str, str] | None = None,
) -> subprocess.CompletedProcess[bytes]:
    safe_env = {
        "PATH": os.environ.get("PATH", ""),
        "HOME": str(cwd),
        "LANG": "C.UTF-8",
        "LC_ALL": "C.UTF-8",
    }
    if env:
        safe_env.update(env)
    return subprocess.run(
        command,
        cwd=cwd,
        env=safe_env,
        capture_output=True,
        timeout=min(timeout, settings.max_timeout_seconds),
        check=False,
    )


def _visual_tool_unavailable(
    settings: RunnerSettings,
    tool_name: str,
) -> dict[str, Any]:
    if settings.office_visual_qa_required:
        raise ToolRejected(f"Runner 缺少 {tool_name}，无法完成办公文件视觉质检")
    return {
        "available": False,
        "tool": tool_name,
        "issues": [],
        "warnings": [f"{tool_name} 不可用，仅完成结构检查"],
    }


def _convert_office_to_pdf(
    target: Path,
    temp_root: Path,
    settings: RunnerSettings,
) -> Path | None:
    output_dir = temp_root / "office-pdf"
    profile_dir = temp_root / "libreoffice-profile"
    output_dir.mkdir(parents=True, exist_ok=True)
    profile_dir.mkdir(parents=True, exist_ok=True)
    try:
        completed = _run_quality_command(
            [
                "soffice",
                "--headless",
                "--nologo",
                "--nodefault",
                "--nolockcheck",
                "--nofirststartwizard",
                f"-env:UserInstallation={profile_dir.resolve().as_uri()}",
                "--convert-to",
                "pdf",
                "--outdir",
                str(output_dir),
                str(target),
            ],
            cwd=temp_root,
            settings=settings,
        )
    except FileNotFoundError:
        _visual_tool_unavailable(settings, "LibreOffice")
        return None
    except subprocess.TimeoutExpired as exc:
        raise ToolRejected("LibreOffice 渲染办公文件超时") from exc
    rendered = output_dir / f"{target.stem}.pdf"
    if completed.returncode != 0 or not rendered.is_file() or rendered.stat().st_size <= 0:
        detail = completed.stderr.decode("utf-8", errors="replace")[:1000]
        raise ToolRejected(f"LibreOffice 无法渲染办公文件：{detail or '未生成 PDF'}")
    return rendered


def _recalculate_xlsx(
    target: Path,
    temp_root: Path,
    settings: RunnerSettings,
) -> Path | None:
    output_dir = temp_root / "recalculated"
    profile_dir = temp_root / "libreoffice-recalc-profile"
    output_dir.mkdir(parents=True, exist_ok=True)
    profile_dir.mkdir(parents=True, exist_ok=True)
    try:
        completed = _run_quality_command(
            [
                "soffice",
                "--headless",
                "--nologo",
                "--nodefault",
                "--nolockcheck",
                "--nofirststartwizard",
                f"-env:UserInstallation={profile_dir.resolve().as_uri()}",
                "--convert-to",
                "xlsx",
                "--outdir",
                str(output_dir),
                str(target),
            ],
            cwd=temp_root,
            settings=settings,
        )
    except (FileNotFoundError, subprocess.TimeoutExpired):
        return None
    rendered = output_dir / f"{target.stem}.xlsx"
    if completed.returncode == 0 and rendered.is_file() and rendered.stat().st_size > 0:
        return rendered
    return None


def _pdf_visual_qa(
    pdf_path: Path,
    temp_root: Path,
    settings: RunnerSettings,
    *,
    text_characters: list[int] | None = None,
    spreadsheet_layout: bool = False,
) -> dict[str, Any]:
    from pypdf import PdfReader

    reader = PdfReader(pdf_path)
    page_count = len(reader.pages)
    if page_count == 0:
        raise ToolRejected("PDF 没有页面")
    rendered_pages = min(page_count, settings.office_visual_max_pages)
    page_texts = [
        (page.extract_text() or "").strip()
        for page in reader.pages[:rendered_pages]
    ]
    if text_characters is None:
        text_characters = [len(text) for text in page_texts]
    render_dir = temp_root / "rendered-pages"
    render_dir.mkdir(parents=True, exist_ok=True)
    prefix = render_dir / "page"
    try:
        completed = _run_quality_command(
            [
                "pdftoppm",
                "-png",
                "-r",
                "110",
                "-f",
                "1",
                "-l",
                str(rendered_pages),
                str(pdf_path),
                str(prefix),
            ],
            cwd=temp_root,
            settings=settings,
        )
    except FileNotFoundError:
        return _visual_tool_unavailable(settings, "Poppler")
    except subprocess.TimeoutExpired as exc:
        raise ToolRejected("PDF 逐页渲染超时") from exc
    if completed.returncode != 0:
        detail = completed.stderr.decode("utf-8", errors="replace")[:1000]
        raise ToolRejected(f"PDF 逐页渲染失败：{detail}")

    from PIL import Image

    images = sorted(render_dir.glob("page-*.png"))
    if len(images) != rendered_pages:
        raise ToolRejected(
            f"PDF 渲染页数不一致：期望 {rendered_pages}，实际 {len(images)}"
        )
    issues: list[dict[str, Any]] = []
    page_stats: list[dict[str, Any]] = []
    for index, image_path in enumerate(images, 1):
        with Image.open(image_path) as image:
            gray = image.convert("L")
            gray.thumbnail((512, 512))
            histogram = gray.histogram()
            total = max(1, sum(histogram))
            white_ratio = sum(histogram[247:]) / total
            dark_ratio = sum(histogram[:12]) / total
            ink_ratio = 1.0 - white_ratio
            glyph_components = _count_glyph_like_components(gray)
            page_stat = {
                "page": index,
                "width": image.width,
                "height": image.height,
                "ink_ratio": round(ink_ratio, 4),
                "dark_ratio": round(dark_ratio, 4),
                "glyph_components": glyph_components,
            }
            page_stats.append(page_stat)
            extracted_chars = (
                text_characters[index - 1]
                if text_characters and index - 1 < len(text_characters)
                else 0
            )
            if ink_ratio < 0.0008 and extracted_chars < 5:
                issues.append({"page": index, "kind": "blank_page"})
            if dark_ratio > 0.85 and extracted_chars < 5:
                issues.append({"page": index, "kind": "nearly_black_page"})
    issues.extend(_pdf_visible_text_issues(page_texts, page_stats))
    if spreadsheet_layout:
        issues.extend(_spreadsheet_pdf_text_issues(page_texts, page_stats))
    warnings: list[str] = []
    if rendered_pages < page_count:
        warnings.append(
            f"文件共 {page_count} 页，视觉质检按上限检查前 {rendered_pages} 页"
        )
    return {
        "available": True,
        "tool": "Poppler",
        "pages": page_count,
        "rendered_pages": rendered_pages,
        "issues": issues,
        "warnings": warnings,
        "page_stats": page_stats,
    }


def _count_glyph_like_components(
    gray: Any,
    *,
    contrast_threshold: int = 36,
) -> int:
    """Count small high-contrast components against the dominant page tone.

    Using a fixed dark-pixel threshold incorrectly treats white text on a dark
    slide as invisible. The median luminance is a stable approximation of the
    page background for both light documents and dark presentation slides.
    """
    width, height = gray.size
    pixels = gray.load()
    histogram = gray.histogram()
    midpoint = max(1, width * height) // 2
    cumulative = 0
    background = 255
    for luminance, count in enumerate(histogram):
        cumulative += count
        if cumulative >= midpoint:
            background = luminance
            break

    def is_foreground(x: int, y: int) -> bool:
        return abs(int(pixels[x, y]) - background) >= contrast_threshold

    seen = bytearray(width * height)
    glyph_components = 0
    for y in range(height):
        for x in range(width):
            offset = y * width + x
            if seen[offset] or not is_foreground(x, y):
                continue
            seen[offset] = 1
            stack = [(x, y)]
            size = 0
            min_x = max_x = x
            min_y = max_y = y
            while stack:
                current_x, current_y = stack.pop()
                size += 1
                min_x = min(min_x, current_x)
                max_x = max(max_x, current_x)
                min_y = min(min_y, current_y)
                max_y = max(max_y, current_y)
                for neighbor_y in range(
                    max(0, current_y - 1), min(height, current_y + 2)
                ):
                    for neighbor_x in range(
                        max(0, current_x - 1), min(width, current_x + 2)
                    ):
                        neighbor_offset = neighbor_y * width + neighbor_x
                        if (
                            seen[neighbor_offset]
                            or not is_foreground(neighbor_x, neighbor_y)
                        ):
                            continue
                        seen[neighbor_offset] = 1
                        stack.append((neighbor_x, neighbor_y))
            component_width = max_x - min_x + 1
            component_height = max_y - min_y + 1
            if (
                2 <= size <= 100
                and component_width <= 20
                and component_height <= 20
            ):
                glyph_components += 1
    return glyph_components


def _pdf_visible_text_issues(
    page_texts: list[str],
    page_stats: list[dict[str, Any]],
) -> list[dict[str, Any]]:
    """Reject PDFs whose text layer exists but whose glyphs are not visibly rendered."""
    issues: list[dict[str, Any]] = []
    for index, page_text in enumerate(page_texts, 1):
        extracted_chars = len(re.sub(r"\s+", "", page_text))
        glyph_components = (
            int(page_stats[index - 1].get("glyph_components") or 0)
            if index - 1 < len(page_stats)
            else 0
        )
        if extracted_chars >= 40 and glyph_components < 10:
            issues.append(
                {
                    "page": index,
                    "kind": "invisible_text_layer",
                    "extracted_characters": extracted_chars,
                    "glyph_components": glyph_components,
                }
            )
    return issues


def _spreadsheet_pdf_text_issues(
    page_texts: list[str],
    page_stats: list[dict[str, Any]],
) -> list[dict[str, Any]]:
    issues: list[dict[str, Any]] = []
    previous_lines: set[str] = set()
    for index, page_text in enumerate(page_texts, 1):
        truncated = re.findall(r"(?<!#)#{3,}(?!#)", page_text)
        if truncated:
            issues.append(
                {
                    "page": index,
                    "kind": "truncated_cell_display",
                    "occurrences": len(truncated),
                }
            )
        current_lines = {
            re.sub(r"\s+", " ", line).strip()
            for line in page_text.splitlines()
            if re.sub(r"\s+", " ", line).strip()
        }
        ink_ratio = (
            float(page_stats[index - 1].get("ink_ratio") or 0)
            if index - 1 < len(page_stats)
            else 0
        )
        if (
            previous_lines
            and len(current_lines) >= 4
            and ink_ratio < 0.02
        ):
            overlap = len(current_lines & previous_lines) / len(current_lines)
            if overlap >= 0.8:
                issues.append(
                    {
                        "page": index,
                        "kind": "repeated_overflow_page",
                        "overlap": round(overlap, 3),
                    }
                )
        previous_lines = current_lines
    return issues


def _document_urls(target: Path) -> int:
    try:
        if target.suffix.lower() in {".docx", ".xlsx", ".xlsm", ".pptx"}:
            with zipfile.ZipFile(target) as archive:
                chunks: list[bytes] = []
                total = 0
                for name in archive.namelist():
                    if not name.endswith((".xml", ".rels")):
                        continue
                    size = archive.getinfo(name).file_size
                    if size > 5_000_000 or total + size > 10_000_000:
                        continue
                    chunks.append(archive.read(name))
                    total += size
                content = b"\n".join(chunks).decode("utf-8", errors="ignore")
            return len(set(URL_PATTERN.findall(content)))
        return len(set(URL_PATTERN.findall(target.read_text(encoding="utf-8", errors="ignore"))))
    except (OSError, zipfile.BadZipFile):
        return 0


def _raise_visual_issues(file_type: str, visual_qa: dict[str, Any]) -> None:
    issues = visual_qa.get("issues") or []
    if issues:
        raise ToolRejected(
            f"{file_type} 视觉质检失败："
            f"{json.dumps(issues[:10], ensure_ascii=False)}"
        )


def _parse_anysearch_payload(action: str, stdout: str) -> dict[str, Any]:
    value = stdout.strip()
    if not value:
        return {}
    try:
        parsed = json.loads(value)
    except json.JSONDecodeError:
        if action == "extract":
            return {"content": value, "characters": len(value)}
        return {"content": value}
    if isinstance(parsed, list):
        return {"items": parsed}
    if isinstance(parsed, dict):
        return parsed
    return {"content": str(parsed)}


def _run_anysearch_command(
    *,
    command: list[str],
    action: str,
    workspace_root: Path,
    env: dict[str, str],
    settings: RunnerSettings,
    extract_url: str | None = None,
) -> ToolResponse:
    cache_key = extract_url.strip() if extract_url else None
    cache_lock = _anysearch_extract_lock(cache_key) if cache_key else None
    if cache_lock is not None:
        cache_lock.acquire()
    try:
        if cache_key:
            cached = _anysearch_extract_cache_get(cache_key)
            if cached is not None:
                cached["cache_hit"] = True
                characters = int(
                    cached.get("characters")
                    or len(str(cached.get("content") or ""))
                )
                return ToolResponse(
                    ok=True,
                    summary=f"AnySearch 已复用网页正文缓存（{characters} 字符）",
                    data=cached,
                )
        try:
            completed = subprocess.run(
                command,
                cwd=workspace_root,
                env=env,
                capture_output=True,
                timeout=min(90, settings.max_timeout_seconds),
                check=False,
            )
        except (FileNotFoundError, subprocess.TimeoutExpired) as exc:
            raise ToolRejected("AnySearch 运行环境不可用或请求超时") from exc
        stdout = completed.stdout[: settings.max_text_bytes].decode(
            "utf-8", errors="replace"
        )
        stderr = completed.stderr[: settings.max_output_bytes].decode(
            "utf-8", errors="replace"
        )
        if completed.returncode != 0:
            return ToolResponse(
                ok=False,
                summary="AnySearch 查询失败",
                data={
                    "returncode": completed.returncode,
                    "stdout": stdout,
                    "stderr": stderr,
                },
            )
        data = _parse_anysearch_payload(action, stdout)
        data["returncode"] = completed.returncode
        data["cache_hit"] = False
        if stderr:
            data["stderr"] = stderr
        if cache_key:
            _anysearch_extract_cache_put(cache_key, data)
        item_count = 0
        for key in ("items", "results"):
            if isinstance(data.get(key), list):
                item_count = len(data[key])
                break
        summary = "AnySearch 查询成功"
        if action == "extract":
            characters = int(
                data.get("characters") or len(str(data.get("content") or ""))
            )
            summary = f"AnySearch 已提取网页正文（{characters} 字符）"
        elif item_count:
            summary = f"AnySearch 已返回 {item_count} 条结果"
        return ToolResponse(ok=True, summary=summary, data=data)
    finally:
        if cache_lock is not None:
            cache_lock.release()


def _skill_path(
    settings: RunnerSettings,
    skill_name: Any,
    raw: Any,
    *,
    must_exist: bool = True,
) -> Path:
    if (
        not isinstance(skill_name, str)
        or not re.fullmatch(r"[a-z0-9][a-z0-9-]{0,63}", skill_name)
    ):
        raise ToolRejected("Skill 名称无效")
    if not isinstance(raw, str) or not raw.strip():
        raise ToolRejected("Skill 路径不能为空")
    pure = PurePath(raw)
    if pure.is_absolute() or ".." in pure.parts:
        raise ToolRejected("Skill 路径无效")
    root = (settings.skills_root / skill_name).resolve(strict=False)
    if not (root / "SKILL.md").is_file():
        raise ToolRejected("Skill 未安装")
    target = (root / pure).resolve(strict=False)
    try:
        target.relative_to(root)
    except ValueError as exc:
        raise ToolRejected("Skill 路径越过边界") from exc
    if must_exist and not target.exists():
        raise ToolRejected("Skill 文件不存在")
    return target


class TaskWorkspace:
    def __init__(self, request: ToolRequest, settings: RunnerSettings | None = None):
        self.settings = settings or get_settings()
        self.root = (
            self.settings.data_root
            / "users"
            / str(request.user_id)
            / "tasks"
            / str(request.task_id)
        ).resolve(strict=False)
        self.root.mkdir(parents=True, exist_ok=True)
        for name in ("input", "workspace", "shared", "output", "logs", "trash"):
            (self.root / name).mkdir(exist_ok=True)

    def path(self, raw: Any, *, must_exist: bool = False) -> Path:
        if not isinstance(raw, str) or not raw.strip():
            raise ToolRejected("路径不能为空")
        pure = PurePath(raw)
        if pure.is_absolute() or ".." in pure.parts:
            raise ToolRejected("路径必须位于任务目录内")
        target = (self.root / pure).resolve(strict=False)
        try:
            target.relative_to(self.root)
        except ValueError as exc:
            raise ToolRejected("路径越过任务目录边界") from exc
        if must_exist and not target.exists():
            raise ToolRejected("文件或目录不存在")
        return target


def _require_approval(request: ToolRequest, message: str) -> None:
    if not request.approval_granted:
        raise ToolRejected(message)


def _arg(arguments: dict[str, Any], key: str) -> Any:
    if key not in arguments:
        raise ToolRejected(f"缺少参数：{key}")
    return arguments[key]


_EXEC_OPEN_PATTERN = re.compile(
    r"""exec\(\s*open\(\s*(['"])([^'"]+)\1\s*\)\.read\(\s*\)\s*\)"""
)


def _python_script_path(workspace: TaskWorkspace, raw: Any) -> Path:
    if not isinstance(raw, str) or not raw.strip():
        raise ToolRejected("script 必须是 workspace 下的 Python 文件路径")
    value = raw.strip()
    wrapper = _EXEC_OPEN_PATTERN.fullmatch(value)
    if wrapper:
        value = wrapper.group(2)
    if "\n" in value or "\r" in value or not value.lower().endswith(".py"):
        raise ToolRejected(
            "script 只能填写文件路径，例如 workspace/create_pptx.py；不能填写 Python 代码"
        )
    if len(PurePath(value).parts) == 1:
        workspace_candidate = f"workspace/{value}"
        if workspace.path(workspace_candidate).is_file():
            value = workspace_candidate
    target = workspace.path(value)
    workspace_dir = workspace.root / "workspace"
    try:
        target.relative_to(workspace_dir)
    except ValueError as exc:
        raise ToolRejected("Python 脚本必须位于 workspace 内") from exc
    if not target.is_file():
        candidates = [
            item.relative_to(workspace.root).as_posix()
            for item in sorted(workspace_dir.rglob("*.py"))
            if item.is_file() and not item.is_symlink()
        ][:10]
        available = "、".join(candidates) if candidates else "无"
        raise ToolRejected(
            "Python 脚本不存在。script 必须直接填写相对路径，例如 "
            f"workspace/create_pptx.py；不能填写代码。当前脚本：{available}"
        )
    return target


def execute(request: ToolRequest, settings: RunnerSettings | None = None) -> ToolResponse:
    settings = settings or get_settings()
    workspace = TaskWorkspace(request, settings)
    try:
        args, agent_scope = enforce_request_scope(request)
    except RunnerScopeError as exc:
        raise ToolRejected(str(exc)) from exc

    if request.tool == "list_files":
        base = workspace.path(args.get("path", "."), must_exist=True)
        if not base.is_dir():
            raise ToolRejected("目标不是目录")
        files = []
        for item in sorted(base.rglob("*")):
            if item.is_symlink():
                continue
            relative = item.relative_to(workspace.root).as_posix()
            files.append({"path": relative, "is_dir": item.is_dir(), "size": item.stat().st_size if item.is_file() else 0})
            if len(files) >= 1000:
                break
        return ToolResponse(ok=True, summary=f"列出 {len(files)} 个项目", data={"items": files})

    if request.tool == "read_text":
        target = workspace.path(_arg(args, "path"), must_exist=True)
        if not target.is_file() or target.stat().st_size > settings.max_text_bytes:
            raise ToolRejected("文件不是可读取的小型文本文件")
        try:
            content = target.read_text(encoding="utf-8")
        except UnicodeDecodeError as exc:
            raise ToolRejected("文件不是 UTF-8 文本，请改用 inspect_document") from exc
        return ToolResponse(ok=True, summary=f"读取 {target.name}", data={"content": content})

    if request.tool == "read_skill_file":
        skill_name = args.get("skill_name", PPT_SKILL_NAME)
        target = _skill_path(settings, skill_name, _arg(args, "path"))
        if not target.is_file() or target.stat().st_size > settings.max_text_bytes:
            raise ToolRejected("Skill 文件不是可读取的小型文本文件")
        try:
            content = target.read_text(encoding="utf-8")
        except UnicodeDecodeError as exc:
            raise ToolRejected("Skill 文件不是 UTF-8 文本，请使用 copy_skill_file") from exc
        skill_root = (settings.skills_root / skill_name).resolve(strict=False)
        return ToolResponse(
            ok=True,
            summary=f"已读取 {skill_name}/{target.relative_to(skill_root).as_posix()}",
            data={"content": content},
        )

    if request.tool == "copy_skill_file":
        skill_name = args.get("skill_name", PPT_SKILL_NAME)
        source = _skill_path(settings, skill_name, _arg(args, "source"))
        target = workspace.path(_arg(args, "target"))
        if not source.is_file():
            raise ToolRejected("Skill 源路径必须是文件")
        if target.exists():
            _require_approval(request, "覆盖已有文件需要用户审批")
        if target.relative_to(workspace.root).parts[0] not in {"workspace", "shared", "output"}:
            raise ToolRejected("Skill 文件只能复制到 workspace、shared 或 output")
        target.parent.mkdir(parents=True, exist_ok=True)
        shutil.copy2(source, target)
        return ToolResponse(
            ok=True,
            summary=f"已复制 Skill 文件 {source.name}",
            data={"path": target.relative_to(workspace.root).as_posix(), "size": target.stat().st_size},
        )

    if request.tool == "validate_swiss_deck":
        target = workspace.path(_arg(args, "path"), must_exist=True)
        if not target.is_file() or target.suffix.lower() != ".html":
            raise ToolRejected("瑞士风校验目标必须是任务目录内的 HTML 文件")
        validator = _skill_path(
            settings, PPT_SKILL_NAME, "scripts/validate-swiss-deck.mjs"
        )
        env = {"PATH": os.environ.get("PATH", ""), "HOME": "/tmp"}
        try:
            completed = subprocess.run(
                ["node", str(validator), str(target)],
                cwd=workspace.root / "workspace",
                env=env,
                capture_output=True,
                timeout=min(60, settings.max_timeout_seconds),
                check=False,
            )
        except FileNotFoundError as exc:
            raise ToolRejected("Runner 未安装 Node.js，无法运行官方校验器") from exc
        except subprocess.TimeoutExpired as exc:
            raise ToolRejected("瑞士风 HTML 校验超时") from exc
        stdout = completed.stdout[: settings.max_output_bytes].decode("utf-8", errors="replace")
        stderr = completed.stderr[: settings.max_output_bytes].decode("utf-8", errors="replace")
        return ToolResponse(
            ok=completed.returncode == 0,
            summary="瑞士风 HTML 校验通过" if completed.returncode == 0 else "瑞士风 HTML 校验失败",
            data={"returncode": completed.returncode, "stdout": stdout, "stderr": stderr},
        )

    if request.tool == "anysearch":
        action = str(_arg(args, "action"))
        if action not in {"search", "batch_search", "extract", "get_sub_domains"}:
            raise ToolRejected("AnySearch action 无效")
        script = _skill_path(
            settings, "anysearch", "scripts/anysearch_cli.js"
        )
        command = ["node", str(script), action]
        if action == "search":
            query = str(_arg(args, "query")).strip()
            if not query or len(query) > 500:
                raise ToolRejected("AnySearch 查询无效")
            command.append(query)
        elif action == "batch_search":
            queries = args.get("queries")
            if (
                not isinstance(queries, list)
                or not queries
                or len(queries) > 5
            ):
                raise ToolRejected("AnySearch 批量查询必须包含 1 到 5 条查询")
            normalized_queries: list[dict[str, Any]] = []
            for item in queries:
                if isinstance(item, str):
                    if not item.strip():
                        raise ToolRejected("AnySearch 批量查询内容不能为空")
                    normalized_queries.append({"query": item.strip()})
                    continue
                if not isinstance(item, dict):
                    raise ToolRejected("AnySearch 批量查询格式无效")
                query = str(item.get("query") or "").strip()
                if not query:
                    raise ToolRejected("AnySearch 批量查询缺少 query")
                normalized = {"query": query}
                for key in ("domain", "sub_domain", "sub_domain_params", "max_results"):
                    if item.get(key) is not None:
                        normalized[key] = item[key]
                normalized_queries.append(normalized)
            command.extend(
                [
                    "--queries",
                    json.dumps(normalized_queries, ensure_ascii=False),
                ]
            )
        elif action == "extract":
            url = str(_arg(args, "url")).strip()
            parsed = urlparse(url)
            if parsed.scheme not in {"http", "https"} or not parsed.netloc:
                raise ToolRejected("AnySearch 提取地址必须是 HTTP/HTTPS URL")
            command.append(url)
        else:
            domains = args.get("domains")
            domain = args.get("domain")
            if isinstance(domains, list) and domains:
                if len(domains) > 10 or any(not isinstance(item, str) for item in domains):
                    raise ToolRejected("AnySearch domains 无效")
                command.extend(["--domains", ",".join(domains)])
            elif isinstance(domain, str) and domain.strip():
                command.extend(["--domain", domain.strip()])
            else:
                raise ToolRejected("AnySearch 需要 domain 或 domains")
        if action in {"search", "batch_search"}:
            if args.get("domain"):
                command.extend(["--domain", str(args["domain"])])
            if args.get("sub_domain"):
                command.extend(["--sub_domain", str(args["sub_domain"])])
            if args.get("sub_domain_params") is not None:
                if not isinstance(args["sub_domain_params"], dict):
                    raise ToolRejected("AnySearch sub_domain_params 必须是对象")
                command.extend(
                    ["--sdp", json.dumps(args["sub_domain_params"], ensure_ascii=False)]
                )
            max_results = int(args.get("max_results", 10))
            if max_results < 1 or max_results > 10:
                raise ToolRejected("AnySearch max_results 必须为 1 到 10")
            command.extend(["--max_results", str(max_results)])
        env = {
            "PATH": os.environ.get("PATH", ""),
            "HOME": "/tmp",
        }
        api_key = os.environ.get("ANYSEARCH_API_KEY", "").strip()
        if api_key:
            env["ANYSEARCH_API_KEY"] = api_key
        return _run_anysearch_command(
            command=command,
            action=action,
            workspace_root=workspace.root / "workspace",
            env=env,
            settings=settings,
            extract_url=url if action == "extract" else None,
        )

    if request.tool == "convert_document":
        source = workspace.path(_arg(args, "source"), must_exist=True)
        target = workspace.path(_arg(args, "target"))
        if not source.is_file() or source.suffix.lower() not in {
            ".docx",
            ".xlsx",
            ".pptx",
        }:
            raise ToolRejected("文档转换来源必须是任务目录内的 DOCX、XLSX 或 PPTX 文件")
        if (
            target.suffix.lower() != ".pdf"
            or target.relative_to(workspace.root).parts[0]
            not in {"workspace", "shared", "output"}
        ):
            raise ToolRejected("文档转换结果必须是 workspace、shared 或 output 中的 PDF")
        if target.exists():
            _require_approval(request, "覆盖已有 PDF 文件需要用户审批")
        with tempfile.TemporaryDirectory(prefix="miniswarm-office-convert-") as raw_temp:
            temp_root = Path(raw_temp)
            output_dir = temp_root / "output"
            profile_dir = temp_root / "profile"
            output_dir.mkdir(parents=True, exist_ok=True)
            profile_dir.mkdir(parents=True, exist_ok=True)
            try:
                completed = _run_quality_command(
                    [
                        "soffice",
                        "--headless",
                        "--nologo",
                        "--nodefault",
                        "--nolockcheck",
                        "--nofirststartwizard",
                        f"-env:UserInstallation={profile_dir.resolve().as_uri()}",
                        "--convert-to",
                        "pdf",
                        "--outdir",
                        str(output_dir),
                        str(source),
                    ],
                    cwd=temp_root,
                    settings=settings,
                )
            except FileNotFoundError as exc:
                raise ToolRejected("Runner 未安装 LibreOffice，无法转换文档") from exc
            except subprocess.TimeoutExpired as exc:
                raise ToolRejected("LibreOffice 文档转换超时") from exc
            converted = output_dir / f"{source.stem}.pdf"
            if (
                completed.returncode != 0
                or not converted.is_file()
                or converted.stat().st_size == 0
            ):
                detail = completed.stderr.decode("utf-8", errors="replace")[:1000]
                raise ToolRejected(f"LibreOffice 文档转换失败：{detail}")
            target.parent.mkdir(parents=True, exist_ok=True)
            shutil.copy2(converted, target)
        return ToolResponse(
            ok=True,
            summary=f"已转换为 {target.name}",
            data={
                "path": target.relative_to(workspace.root).as_posix(),
                "size": target.stat().st_size,
                "source": source.relative_to(workspace.root).as_posix(),
            },
        )

    if request.tool == "convert_to_markdown":
        source = workspace.path(_arg(args, "source"), must_exist=True)
        target = workspace.path(_arg(args, "target"))
        if not source.is_file():
            raise ToolRejected("MarkItDown 来源必须是文件")
        if (
            target.suffix.lower() != ".md"
            or target.relative_to(workspace.root).parts[0]
            not in {"workspace", "shared", "output"}
        ):
            raise ToolRejected("Markdown 只能写入 workspace、shared 或 output")
        if target.exists():
            _require_approval(request, "覆盖已有 Markdown 文件需要用户审批")
        try:
            from markitdown import MarkItDown

            result = MarkItDown(enable_plugins=False).convert(str(source))
            content = result.text_content
        except Exception as exc:
            raise ToolRejected(f"MarkItDown 转换失败：{type(exc).__name__}") from exc
        if len(content.encode("utf-8")) > settings.max_text_bytes:
            raise ToolRejected("MarkItDown 输出超过文本大小限制")
        target.parent.mkdir(parents=True, exist_ok=True)
        target.write_text(content, encoding="utf-8")
        return ToolResponse(
            ok=True,
            summary=f"已转换为 {target.name}",
            data={
                "path": target.relative_to(workspace.root).as_posix(),
                "size": target.stat().st_size,
            },
        )

    if request.tool == "inspect_document":
        target = workspace.path(_arg(args, "path"), must_exist=True)
        if not target.is_file():
            raise ToolRejected("检查目标必须是文件")
        if target.stat().st_size > settings.max_office_bytes:
            raise ToolRejected("文件超过办公质检大小限制")
        suffix = target.suffix.lower()
        summary_detail = ""
        try:
            if suffix == ".docx":
                from docx import Document
                from docx.enum.table import WD_ROW_HEIGHT_RULE

                document = Document(target)
                nonempty = sum(
                    1 for paragraph in document.paragraphs if paragraph.text.strip()
                )
                table_cells = sum(
                    1
                    for table in document.tables
                    for row in table.rows
                    for cell in row.cells
                    if cell.text.strip()
                )
                if nonempty == 0 and table_cells == 0:
                    raise ToolRejected("DOCX 没有有效正文或表格内容")
                headings = sum(
                    1
                    for paragraph in document.paragraphs
                    if paragraph.text.strip()
                    and (paragraph.style.name or "").casefold().startswith("heading")
                )
                numbered_paragraphs = sum(
                    1
                    for paragraph in document.paragraphs
                    if paragraph._p.pPr is not None
                    and paragraph._p.pPr.numPr is not None
                )
                fake_list_items = sum(
                    1
                    for paragraph in document.paragraphs
                    if re.match(
                        r"^\s*(?:[•●▪◦\-]|\d+[.)、])\s+\S",
                        paragraph.text,
                    )
                    and not (
                        paragraph._p.pPr is not None
                        and paragraph._p.pPr.numPr is not None
                    )
                )
                fixed_height_rows = sum(
                    1
                    for table in document.tables
                    for row in table.rows
                    if row.height_rule == WD_ROW_HEIGHT_RULE.EXACTLY
                    and any(cell.text.strip() for cell in row.cells)
                )
                explicit_font_sizes = [
                    round(run.font.size.pt, 2)
                    for paragraph in document.paragraphs
                    for run in paragraph.runs
                    if run.text.strip() and run.font.size is not None
                ]
                warnings: list[str] = []
                if fake_list_items:
                    warnings.append(f"发现 {fake_list_items} 个疑似手工列表项")
                if fixed_height_rows:
                    warnings.append(
                        f"发现 {fixed_height_rows} 个固定高度表格行，需留意文字裁切"
                    )
                with tempfile.TemporaryDirectory(prefix="miniswarm-docx-qa-") as raw_temp:
                    temp_root = Path(raw_temp)
                    rendered_pdf = _convert_office_to_pdf(
                        target, temp_root, settings
                    )
                    if rendered_pdf is None:
                        visual_qa = _visual_tool_unavailable(
                            settings, "LibreOffice"
                        )
                    else:
                        from pypdf import PdfReader

                        rendered_reader = PdfReader(rendered_pdf)
                        text_characters = [
                            len((page.extract_text() or "").strip())
                            for page in rendered_reader.pages
                        ]
                        visual_qa = _pdf_visual_qa(
                            rendered_pdf,
                            temp_root,
                            settings,
                            text_characters=text_characters,
                        )
                _raise_visual_issues("DOCX", visual_qa)
                data = {
                    "type": "docx",
                    "paragraphs": len(document.paragraphs),
                    "nonempty_paragraphs": nonempty,
                    "headings": headings,
                    "numbered_paragraphs": numbered_paragraphs,
                    "tables": len(document.tables),
                    "nonempty_table_cells": table_cells,
                    "sections": len(document.sections),
                    "inline_images": len(document.inline_shapes),
                    "fake_list_items": fake_list_items,
                    "fixed_height_rows": fixed_height_rows,
                    "minimum_explicit_font_pt": (
                        min(explicit_font_sizes) if explicit_font_sizes else None
                    ),
                    "urls": _document_urls(target),
                    "warnings": warnings,
                    "visual_qa": visual_qa,
                }
                summary_detail = (
                    f"{visual_qa.get('pages', '?')} 页，"
                    f"{len(document.tables)} 个表格，视觉质检通过"
                    if visual_qa.get("available")
                    else f"{len(document.tables)} 个表格，结构检查通过"
                )
            elif suffix in {".xlsx", ".xlsm"}:
                from openpyxl import load_workbook

                keep_vba = suffix == ".xlsm"
                workbook = load_workbook(
                    target,
                    read_only=False,
                    data_only=False,
                    keep_links=False,
                    keep_vba=keep_vba,
                )
                formula_count = 0
                nonempty_cells = 0
                formula_reference_errors: list[str] = []
                numeric_general_cells = 0
                scanned_cells = 0
                scan_limit = 250_000
                sheets: list[dict[str, Any]] = []
                for sheet in workbook.worksheets:
                    sheet_nonempty = 0
                    sheet_formulas = 0
                    for row in sheet.iter_rows():
                        for cell in row:
                            if scanned_cells >= scan_limit:
                                break
                            scanned_cells += 1
                            if cell.value is None:
                                continue
                            nonempty_cells += 1
                            sheet_nonempty += 1
                            if isinstance(cell.value, str) and cell.value.startswith("="):
                                formula_count += 1
                                sheet_formulas += 1
                                if "#REF!" in cell.value.upper():
                                    formula_reference_errors.append(
                                        f"{sheet.title}!{cell.coordinate}"
                                    )
                            if (
                                isinstance(cell.value, (int, float))
                                and cell.number_format == "General"
                            ):
                                numeric_general_cells += 1
                        if scanned_cells >= scan_limit:
                            break
                    sheets.append(
                        {
                            "name": sheet.title,
                            "state": sheet.sheet_state,
                            "rows": sheet.max_row,
                            "columns": sheet.max_column,
                            "nonempty_cells": sheet_nonempty,
                            "formulas": sheet_formulas,
                            "tables": len(sheet.tables),
                            "charts": len(sheet._charts),
                            "merged_ranges": len(sheet.merged_cells.ranges),
                            "freeze_panes": (
                                str(sheet.freeze_panes)
                                if sheet.freeze_panes is not None
                                else None
                            ),
                            "auto_filter": sheet.auto_filter.ref,
                        }
                    )
                    if scanned_cells >= scan_limit:
                        break
                if nonempty_cells == 0:
                    workbook.close()
                    raise ToolRejected("XLSX 没有有效单元格内容")
                with tempfile.TemporaryDirectory(prefix="miniswarm-xlsx-qa-") as raw_temp:
                    temp_root = Path(raw_temp)
                    recalculated = _recalculate_xlsx(
                        target, temp_root, settings
                    )
                    value_source = recalculated or target
                    value_book = load_workbook(
                        value_source,
                        read_only=True,
                        data_only=True,
                        keep_links=False,
                    )
                    formula_result_errors: list[str] = []
                    for sheet in value_book.worksheets:
                        for row in sheet.iter_rows():
                            for cell in row:
                                if (
                                    isinstance(cell.value, str)
                                    and cell.value.upper()
                                    in {
                                        "#REF!",
                                        "#DIV/0!",
                                        "#VALUE!",
                                        "#NAME?",
                                        "#N/A",
                                        "#NUM!",
                                        "#NULL!",
                                    }
                                ):
                                    formula_result_errors.append(
                                        f"{sheet.title}!{cell.coordinate}={cell.value}"
                                    )
                                    if len(formula_result_errors) >= 50:
                                        break
                            if len(formula_result_errors) >= 50:
                                break
                        if len(formula_result_errors) >= 50:
                            break
                    value_book.close()
                    rendered_pdf = _convert_office_to_pdf(
                        target, temp_root, settings
                    )
                    if rendered_pdf is None:
                        visual_qa = _visual_tool_unavailable(
                            settings, "LibreOffice"
                        )
                    else:
                        visual_qa = _pdf_visual_qa(
                            rendered_pdf,
                            temp_root,
                            settings,
                            spreadsheet_layout=True,
                        )
                workbook.close()
                if formula_reference_errors or formula_result_errors:
                    preview = [
                        *formula_reference_errors[:10],
                        *formula_result_errors[:10],
                    ]
                    raise ToolRejected(
                        "XLSX 检测到公式错误："
                        + "、".join(preview)
                    )
                _raise_visual_issues("XLSX", visual_qa)
                warnings = []
                if scanned_cells >= scan_limit:
                    warnings.append(
                        f"工作簿较大，结构扫描按上限检查前 {scan_limit} 个单元格"
                    )
                if numeric_general_cells:
                    warnings.append(
                        f"{numeric_general_cells} 个数值单元格仍使用 General 格式"
                    )
                data = {
                    "type": "xlsx",
                    "sheets": sheets,
                    "nonempty_cells": nonempty_cells,
                    "formula_count": formula_count,
                    "formula_errors": [],
                    "recalculated_with_libreoffice": recalculated is not None,
                    "numeric_general_cells": numeric_general_cells,
                    "urls": _document_urls(target),
                    "warnings": warnings,
                    "visual_qa": visual_qa,
                }
                summary_detail = (
                    f"{len(sheets)} 个工作表，{formula_count} 个公式，"
                    f"{visual_qa.get('pages', '?')} 个渲染页，检查通过"
                    if visual_qa.get("available")
                    else f"{len(sheets)} 个工作表，{formula_count} 个公式，结构检查通过"
                )
            elif suffix == ".pptx":
                from pptx import Presentation

                presentation = Presentation(target)
                text_shapes = sum(
                    1 for slide in presentation.slides for shape in slide.shapes
                    if hasattr(shape, "text") and shape.text.strip()
                )
                layout_issues = _pptx_layout_issues(presentation)
                data = {
                    "type": "pptx",
                    "slides": len(presentation.slides),
                    "text_shapes": text_shapes,
                    "layout_issues": layout_issues,
                }
                if layout_issues:
                    preview = json.dumps(layout_issues[:5], ensure_ascii=False)
                    raise ToolRejected(f"PPTX 检测到文本越界或重叠：{preview}")
                summary_detail = f"{len(presentation.slides)} 页，布局检查通过"
            elif suffix == ".pdf":
                from pypdf import PdfReader

                reader = PdfReader(target)
                if reader.is_encrypted:
                    raise ToolRejected("PDF 已加密，无法完成内容与视觉质检")
                if not reader.pages:
                    raise ToolRejected("PDF 没有页面")
                text_characters = [
                    len((page.extract_text() or "").strip())
                    for page in reader.pages
                ]
                page_sizes = [
                    [
                        round(float(page.mediabox.width), 2),
                        round(float(page.mediabox.height), 2),
                    ]
                    for page in reader.pages
                ]
                with tempfile.TemporaryDirectory(prefix="miniswarm-pdf-qa-") as raw_temp:
                    visual_qa = _pdf_visual_qa(
                        target,
                        Path(raw_temp),
                        settings,
                        text_characters=text_characters,
                    )
                _raise_visual_issues("PDF", visual_qa)
                data = {
                    "type": "pdf",
                    "pages": len(reader.pages),
                    "encrypted": False,
                    "text_characters": sum(text_characters),
                    "pages_without_text": sum(
                        1 for count in text_characters if count == 0
                    ),
                    "page_sizes": page_sizes[:20],
                    "urls": len(
                        set(
                            URL_PATTERN.findall(
                                "\n".join(
                                    page.extract_text() or ""
                                    for page in reader.pages
                                )
                            )
                        )
                    ),
                    "visual_qa": visual_qa,
                }
                summary_detail = (
                    f"{len(reader.pages)} 页，逐页渲染检查通过"
                    if visual_qa.get("available")
                    else f"{len(reader.pages)} 页，结构检查通过"
                )
            elif suffix in {".png", ".jpg", ".jpeg", ".webp", ".gif", ".bmp", ".tiff"}:
                from PIL import Image

                with Image.open(target) as image:
                    image.verify()
                with Image.open(target) as image:
                    data = {"type": "image", "format": image.format, "width": image.width, "height": image.height, "mode": image.mode}
            elif suffix in {".html", ".htm"}:
                if target.stat().st_size > settings.max_text_bytes:
                    raise ToolRejected("HTML 超过检查大小限制")
                content = target.read_text(encoding="utf-8")
                parser = _HTMLSummaryParser()
                parser.feed(content)
                parser.close()
                data = {
                    "type": "html",
                    "title": " ".join(parser.title_parts)[:300],
                    "elements": parser.elements,
                    "text_chars": parser.text_chars,
                    "scripts": parser.scripts,
                    "forms": parser.forms,
                    "slides": parser.slides,
                }
                if parser.elements == 0 or parser.text_chars == 0:
                    raise ToolRejected("HTML 缺少有效结构或可见文本")
            elif suffix == ".csv":
                rows = 0
                max_columns = 0
                with target.open("r", encoding="utf-8-sig", newline="") as source:
                    for row in csv.reader(source):
                        rows += 1
                        max_columns = max(max_columns, len(row))
                        if rows >= 10_000:
                            break
                if rows == 0 or max_columns == 0:
                    raise ToolRejected("CSV 没有有效数据")
                data = {"type": "csv", "rows_scanned": rows, "max_columns": max_columns}
            elif suffix == ".zip":
                with zipfile.ZipFile(target) as archive:
                    members = archive.infolist()
                    broken = archive.testzip()
                    if broken:
                        raise ToolRejected(f"ZIP 包含损坏文件：{broken}")
                    if not members:
                        raise ToolRejected("ZIP 为空")
                    data = {
                        "type": "zip",
                        "entries": len(members),
                        "uncompressed_size": sum(item.file_size for item in members),
                    }
            elif suffix in {
                ".txt", ".md", ".json", ".xml", ".yaml", ".yml",
                ".py", ".js", ".ts", ".css", ".svg",
            }:
                if target.stat().st_size > settings.max_text_bytes:
                    raise ToolRejected("文本文件超过检查大小限制")
                content = target.read_text(encoding="utf-8")
                if not content.strip():
                    raise ToolRejected("文本文件为空")
                data = {
                    "type": "text",
                    "characters": len(content),
                    "lines": content.count("\n") + 1,
                    "nonempty_lines": sum(1 for line in content.splitlines() if line.strip()),
                }
                if suffix == ".json":
                    parsed = json.loads(content)
                    data["json_type"] = type(parsed).__name__
            else:
                raise ToolRejected("不支持检查此文件类型")
        except ToolRejected:
            raise
        except Exception as exc:
            raise ToolRejected(f"文档无法打开或已损坏：{type(exc).__name__}") from exc
        data["size"] = target.stat().st_size
        summary = f"{target.name} 检查通过"
        if summary_detail:
            summary = f"{summary}（{summary_detail}）"
        return ToolResponse(ok=True, summary=summary, data=data)

    if request.tool == "write_text":
        target = workspace.path(_arg(args, "path"))
        content = _arg(args, "content")
        if not isinstance(content, str) or len(content.encode("utf-8")) > settings.max_text_bytes:
            raise ToolRejected("文本内容无效或过大")
        if target.exists():
            _require_approval(request, "覆盖已有文件需要用户审批")
        if not target.relative_to(workspace.root).parts[0] in {"workspace", "shared", "output"}:
            raise ToolRejected("只能写入 workspace、shared 或 output")
        target.parent.mkdir(parents=True, exist_ok=True)
        target.write_text(content, encoding="utf-8")
        return ToolResponse(ok=True, summary=f"已写入 {target.name}", data={"size": target.stat().st_size})

    if request.tool in {"copy_file", "move_file"}:
        source = workspace.path(_arg(args, "source"), must_exist=True)
        target = workspace.path(_arg(args, "target"))
        if not source.is_file():
            raise ToolRejected("源路径必须是文件")
        if target.exists():
            _require_approval(request, "覆盖已有文件需要用户审批")
        target.parent.mkdir(parents=True, exist_ok=True)
        if request.tool == "copy_file":
            shutil.copy2(source, target)
            action = "复制"
        else:
            _require_approval(request, "移动已有文件需要用户审批")
            shutil.move(str(source), str(target))
            action = "移动"
        return ToolResponse(ok=True, summary=f"已{action}文件", data={"path": target.relative_to(workspace.root).as_posix()})

    if request.tool == "move_to_trash":
        _require_approval(request, "删除文件需要用户审批")
        source = workspace.path(_arg(args, "path"), must_exist=True)
        if source == workspace.root or source.parent == workspace.root:
            raise ToolRejected("不能删除任务根目录或系统目录")
        destination = workspace.root / "trash" / source.relative_to(workspace.root)
        destination.parent.mkdir(parents=True, exist_ok=True)
        if destination.exists():
            destination = destination.with_name(f"{destination.stem}-{request.request_id}{destination.suffix}")
        shutil.move(str(source), str(destination))
        return ToolResponse(ok=True, summary="已移入回收站", data={"path": destination.relative_to(workspace.root).as_posix()})

    if request.tool == "create_directory":
        target = workspace.path(_arg(args, "path"))
        if not target.relative_to(workspace.root).parts[0] in {"workspace", "shared", "output"}:
            raise ToolRejected("只能在工作目录中创建文件夹")
        target.mkdir(parents=True, exist_ok=True)
        return ToolResponse(ok=True, summary="已创建文件夹")

    if request.tool == "create_zip":
        sources = _arg(args, "sources")
        if not isinstance(sources, list) or not sources or len(sources) > 100:
            raise ToolRejected("压缩来源必须是 1 到 100 个路径")
        target = workspace.path(_arg(args, "target"))
        if target.suffix.lower() != ".zip" or target.relative_to(workspace.root).parts[0] != "output":
            raise ToolRejected("ZIP 只能写入 output 目录")
        if target.exists():
            _require_approval(request, "覆盖已有 ZIP 需要用户审批")
        target.parent.mkdir(parents=True, exist_ok=True)
        with zipfile.ZipFile(target, "w", compression=zipfile.ZIP_DEFLATED) as archive:
            for raw in sources:
                source = workspace.path(raw, must_exist=True)
                candidates = [source] if source.is_file() else [item for item in source.rglob("*") if item.is_file() and not item.is_symlink()]
                for item in candidates:
                    relative = item.relative_to(workspace.root)
                    if relative.parts[0] != "trash":
                        archive.write(item, relative.as_posix())
        return ToolResponse(ok=True, summary="ZIP 已创建", data={"path": target.relative_to(workspace.root).as_posix(), "size": target.stat().st_size})

    if request.tool == "run_python":
        script = _python_script_path(workspace, _arg(args, "script"))
        timeout = int(args.get("timeout_seconds", 60))
        if timeout < 1 or timeout > settings.max_timeout_seconds:
            raise ToolRejected("执行超时参数越界")
        env = {
            "PATH": os.environ.get("PATH", ""),
            "PYTHONIOENCODING": "utf-8",
            "PYTHONDONTWRITEBYTECODE": "1",
            "PYTEST_DISABLE_PLUGIN_AUTOLOAD": "1",
            "HOME": "/tmp",
        }
        command = [sys.executable, "-I", str(script)]
        process_cwd = workspace.root / "workspace"
        if agent_scope is not None:
            runtime_home = workspace.root / agent_scope.workspace / ".runtime"
            runtime_home.mkdir(parents=True, exist_ok=True)
            env.update({
                "HOME": str(runtime_home),
                "TMPDIR": str(runtime_home),
                "XDG_CACHE_HOME": str(runtime_home / ".cache"),
                "MPLCONFIGDIR": str(runtime_home / ".matplotlib"),
                "MINISWARM_TASK_ROOT": str(workspace.root),
                "MINISWARM_AGENT_WORKSPACE": agent_scope.workspace,
                "MINISWARM_AGENT_OUTPUT": agent_scope.output,
            })
            # Signed Agent paths are task-root-relative. Running scoped scripts
            # from the task root makes the documented paths resolve exactly once
            # instead of producing workspace/workspace/... false violations.
            process_cwd = workspace.root
            command = [
                sys.executable,
                "-I",
                str(Path(__file__).with_name("sandbox_runner.py")),
                "--task-root",
                str(workspace.root),
                "--script",
                str(script),
                "--read-roots",
                json.dumps(agent_scope.readable_roots),
                "--write-roots",
                json.dumps(agent_scope.writable_roots),
                "--role",
                agent_scope.role,
                "--output-root",
                agent_scope.output,
            ]
        try:
            completed = subprocess.run(
                command,
                cwd=process_cwd,
                env=env,
                capture_output=True,
                timeout=timeout,
                check=False,
            )
        except subprocess.TimeoutExpired as exc:
            raise ToolRejected("Python 执行超时") from exc
        stdout = completed.stdout[: settings.max_output_bytes].decode("utf-8", errors="replace")
        stderr = completed.stderr[: settings.max_output_bytes].decode("utf-8", errors="replace")
        return ToolResponse(
            ok=completed.returncode == 0,
            summary="Python 执行成功" if completed.returncode == 0 else "Python 执行失败",
            data={"returncode": completed.returncode, "stdout": stdout, "stderr": stderr},
        )

    if request.tool == "run_tests":
        target = workspace.path(args.get("path", "workspace"), must_exist=True)
        workspace_dir = workspace.root / "workspace"
        try:
            target.relative_to(workspace_dir)
        except ValueError as exc:
            raise ToolRejected("测试路径必须位于 workspace 内") from exc
        timeout = int(args.get("timeout_seconds", 120))
        if timeout < 1 or timeout > settings.max_timeout_seconds:
            raise ToolRejected("测试超时参数越界")
        env = {
            "PATH": os.environ.get("PATH", ""),
            "PYTHONIOENCODING": "utf-8",
            "PYTHONDONTWRITEBYTECODE": "1",
            "PYTEST_DISABLE_PLUGIN_AUTOLOAD": "1",
            "HOME": "/tmp",
        }
        command = [sys.executable, "-I", "-m", "pytest", "-q", str(target)]
        process_cwd = workspace_dir
        if agent_scope is not None:
            runtime_home = workspace.root / agent_scope.workspace / ".runtime"
            runtime_home.mkdir(parents=True, exist_ok=True)
            env.update({
                "HOME": str(runtime_home),
                "TMPDIR": str(runtime_home),
                "XDG_CACHE_HOME": str(runtime_home / ".cache"),
                "MPLCONFIGDIR": str(runtime_home / ".matplotlib"),
                "MINISWARM_TASK_ROOT": str(workspace.root),
                "MINISWARM_AGENT_WORKSPACE": agent_scope.workspace,
                "MINISWARM_AGENT_OUTPUT": agent_scope.output,
            })
            process_cwd = workspace.root
            command = [
                sys.executable,
                "-I",
                str(Path(__file__).with_name("sandbox_runner.py")),
                "--task-root",
                str(workspace.root),
                "--pytest-target",
                str(target),
                "--read-roots",
                json.dumps(agent_scope.readable_roots),
                "--write-roots",
                json.dumps(agent_scope.writable_roots),
                "--role",
                agent_scope.role,
                "--output-root",
                agent_scope.output,
            ]
        try:
            completed = subprocess.run(
                command,
                cwd=process_cwd,
                env=env,
                capture_output=True,
                timeout=timeout,
                check=False,
            )
        except subprocess.TimeoutExpired as exc:
            raise ToolRejected("测试执行超时") from exc
        stdout = completed.stdout[: settings.max_output_bytes].decode("utf-8", errors="replace")
        stderr = completed.stderr[: settings.max_output_bytes].decode("utf-8", errors="replace")
        return ToolResponse(
            ok=completed.returncode == 0,
            summary="测试通过" if completed.returncode == 0 else "测试失败",
            data={"returncode": completed.returncode, "stdout": stdout, "stderr": stderr},
        )

    raise ToolRejected("不支持的工具")
