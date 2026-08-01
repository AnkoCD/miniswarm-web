from uuid import uuid4
import json
import subprocess
import sys
from pathlib import Path
from types import SimpleNamespace

import pytest

from runner_app.schemas import ToolRequest
from runner_app.tools import (
    ToolRejected,
    _count_glyph_like_components,
    _pdf_visible_text_issues,
    _spreadsheet_pdf_text_issues,
    execute,
)


def request(tool, arguments, *, approved=False):
    return ToolRequest(
        request_id=uuid4(), user_id=uuid4(), task_id=uuid4(), tool=tool,
        arguments=arguments, approval_granted=approved,
    )


def test_rejects_path_traversal(settings):
    with pytest.raises(ToolRejected, match="任务目录"):
        execute(request("read_text", {"path": "../../secret.txt"}), settings)


def test_write_and_read_text(settings):
    item = request("write_text", {"path": "output/result.txt", "content": "hello"})
    assert execute(item, settings).ok
    read = item.model_copy(update={"tool": "read_text", "arguments": {"path": "output/result.txt"}})
    assert execute(read, settings).data["content"] == "hello"


def test_read_text_rejects_binary_as_tool_error(settings):
    item = request("read_text", {"path": "output/binary.pptx"})
    root = Path(settings.data_root) / "users" / str(item.user_id) / "tasks" / str(item.task_id)
    output = root / "output"
    output.mkdir(parents=True, exist_ok=True)
    (output / "binary.pptx").write_bytes(b"\x80\x81\x82\xff")
    with pytest.raises(ToolRejected, match="inspect_document"):
        execute(item, settings)


def test_read_and_copy_installed_ppt_skill(settings):
    read = request("read_skill_file", {"path": "SKILL.md"})
    assert "Test PPT Skill" in execute(read, settings).data["content"]
    copied = execute(
        read.model_copy(
            update={
                "tool": "copy_skill_file",
                "arguments": {
                    "source": "assets/template.html",
                    "target": "workspace/ppt/index.html",
                },
            }
        ),
        settings,
    )
    assert copied.ok
    assert (settings.data_root / "users" / str(read.user_id) / "tasks" / str(read.task_id) / "workspace" / "ppt" / "index.html").is_file()


def test_read_named_installed_skill(settings):
    result = execute(
        request(
            "read_skill_file",
            {"skill_name": "anysearch", "path": "SKILL.md"},
        ),
        settings,
    )
    assert "AnySearch" in result.data["content"]


def test_skill_path_traversal_is_rejected(settings):
    with pytest.raises(ToolRejected, match="Skill 路径"):
        execute(request("read_skill_file", {"path": "../secret"}), settings)


def test_validate_swiss_deck_uses_fixed_official_script(settings, monkeypatch):
    item = request("write_text", {"path": "output/deck.html", "content": '<section class="slide"></section>'})
    execute(item, settings)
    calls = []

    def fake_run(command, **kwargs):
        calls.append(command)
        return subprocess.CompletedProcess(command, 0, stdout=b"passed", stderr=b"")

    monkeypatch.setattr("runner_app.tools.subprocess.run", fake_run)
    result = execute(
        item.model_copy(update={"tool": "validate_swiss_deck", "arguments": {"path": "output/deck.html"}}),
        settings,
    )
    assert result.ok
    assert calls[0][0] == "node"
    assert calls[0][1].endswith("validate-swiss-deck.mjs")


def test_anysearch_uses_fixed_cli_and_bounded_arguments(settings, monkeypatch):
    calls = []

    def fake_run(command, **kwargs):
        calls.append((command, kwargs))
        return subprocess.CompletedProcess(command, 0, stdout=b'{"items":[]}', stderr=b"")

    monkeypatch.setattr("runner_app.tools.subprocess.run", fake_run)
    result = execute(
        request(
            "anysearch",
            {"action": "search", "query": "MiniSwarm", "max_results": 3},
        ),
        settings,
    )
    assert result.ok
    assert calls[0][0][-4:] == ["search", "MiniSwarm", "--max_results", "3"]
    assert "shell" not in calls[0][1]


def test_anysearch_batch_uses_official_queries_json_and_parses_results(settings, monkeypatch):
    calls = []

    def fake_run(command, **kwargs):
        calls.append(command)
        payload = b'{"results":[{"title":"Official","url":"https://example.com/source","snippet":"Primary source"}]}'
        return subprocess.CompletedProcess(command, 0, stdout=payload, stderr=b"")

    monkeypatch.setattr("runner_app.tools.subprocess.run", fake_run)
    result = execute(
        request(
            "anysearch",
            {
                "action": "batch_search",
                "queries": [
                    {"query": "official specification", "domain": "code"},
                    "independent comparison",
                ],
                "max_results": 4,
            },
        ),
        settings,
    )
    assert result.ok
    assert calls[0][-2:] == ["--max_results", "4"]
    assert "--queries" in calls[0]
    encoded = calls[0][calls[0].index("--queries") + 1]
    assert json.loads(encoded)[0]["domain"] == "code"
    assert result.data["results"][0]["url"] == "https://example.com/source"


def test_anysearch_extract_reuses_successful_content_cache(settings, monkeypatch):
    calls = []

    def fake_run(command, **kwargs):
        calls.append(command)
        payload = b'{"content":"verified source body","characters":20}'
        return subprocess.CompletedProcess(command, 0, stdout=payload, stderr=b"")

    monkeypatch.setattr("runner_app.tools.subprocess.run", fake_run)
    monkeypatch.setattr("runner_app.tools._ANYSEARCH_EXTRACT_CACHE", {})
    first = execute(
        request(
            "anysearch",
            {"action": "extract", "url": "https://example.com/cache-test"},
        ),
        settings,
    )
    second = execute(
        request(
            "anysearch",
            {"action": "extract", "url": "https://example.com/cache-test"},
        ),
        settings,
    )
    assert first.ok and second.ok
    assert len(calls) == 1
    assert first.data["cache_hit"] is False
    assert second.data["cache_hit"] is True
    assert "缓存" in second.summary


def test_spreadsheet_visual_text_gate_detects_truncation_and_overflow_page():
    issues = _spreadsheet_pdf_text_issues(
        [
            "仪表板\n总预算 总实际 差异\n### ### ¥2,617.79\n"
            "01 月 02 月 03 月\n趋势图\n预算合计\n实际合计",
            "01 月 02 月 03 月\n趋势图\n预算合计\n实际合计",
        ],
        [
            {"page": 1, "ink_ratio": 0.12},
            {"page": 2, "ink_ratio": 0.009},
        ],
    )
    assert {item["kind"] for item in issues} == {
        "truncated_cell_display",
        "repeated_overflow_page",
    }


def test_spreadsheet_visual_text_gate_allows_sparse_unique_cover():
    assert _spreadsheet_pdf_text_issues(
        ["预算工作簿\n参数说明", "月度数据\n一月\n二月\n三月"],
        [
            {"page": 1, "ink_ratio": 0.008},
            {"page": 2, "ink_ratio": 0.15},
        ],
    ) == []


def test_pdf_visual_text_gate_detects_invisible_text_layer():
    assert _pdf_visible_text_issues(
        ["办公操作手册适用范围全体员工版本一" * 3],
        [{"page": 1, "ink_ratio": 0.01, "glyph_components": 2}],
    ) == [
        {
            "page": 1,
            "kind": "invisible_text_layer",
            "extracted_characters": 51,
            "glyph_components": 2,
        }
    ]


def test_pdf_visual_text_gate_allows_visible_sparse_cover():
    assert _pdf_visible_text_issues(
        ["办公操作手册适用范围全体员工版本一" * 3],
        [{"page": 1, "ink_ratio": 0.01, "glyph_components": 24}],
    ) == []


def test_glyph_component_detector_supports_light_and_dark_pages():
    from PIL import Image, ImageDraw

    def sample(background, foreground):
        image = Image.new("L", (160, 80), background)
        draw = ImageDraw.Draw(image)
        for index in range(12):
            x = 5 + (index % 6) * 24
            y = 8 + (index // 6) * 28
            draw.rectangle((x, y, x + 4, y + 7), fill=foreground)
        return image

    assert _count_glyph_like_components(sample(255, 20)) >= 10
    assert _count_glyph_like_components(sample(20, 245)) >= 10
    assert _count_glyph_like_components(Image.new("L", (160, 80), 20)) == 0


def test_convert_document_uses_bounded_libreoffice_conversion(settings, monkeypatch):
    source = request(
        "write_text", {"path": "workspace/source.docx", "content": "placeholder"}
    )
    execute(source, settings)

    def fake_quality_command(command, *, cwd, settings):
        output_dir = Path(command[command.index("--outdir") + 1])
        output_dir.mkdir(parents=True, exist_ok=True)
        (output_dir / "source.pdf").write_bytes(b"%PDF-1.4\nfixture\n")
        return subprocess.CompletedProcess(command, 0, b"converted", b"")

    monkeypatch.setattr(
        "runner_app.tools._run_quality_command",
        fake_quality_command,
    )
    result = execute(
        ToolRequest.model_validate(
            {
                **source.model_dump(),
                "tool": "convert_document",
                "arguments": {
                    "source": "workspace/source.docx",
                    "target": "output/source.pdf",
                },
            }
        ),
        settings,
    )
    assert result.ok
    assert result.data["path"] == "output/source.pdf"


def test_markitdown_writes_only_task_markdown(settings, monkeypatch):
    class FakeMarkItDown:
        def __init__(self, **kwargs):
            assert kwargs == {"enable_plugins": False}

        def convert(self, path):
            assert path.endswith("source.txt")
            return SimpleNamespace(text_content="# Converted")

    monkeypatch.setitem(
        sys.modules, "markitdown", SimpleNamespace(MarkItDown=FakeMarkItDown)
    )
    source = request(
        "write_text", {"path": "workspace/source.txt", "content": "source"}
    )
    execute(source, settings)
    result = execute(
        source.model_copy(
            update={
                "tool": "convert_to_markdown",
                "arguments": {
                    "source": "workspace/source.txt",
                    "target": "output/source.md",
                },
            }
        ),
        settings,
    )
    assert result.ok
    target = (
        settings.data_root
        / "users"
        / str(source.user_id)
        / "tasks"
        / str(source.task_id)
        / "output"
        / "source.md"
    )
    assert target.read_text(encoding="utf-8") == "# Converted"


def test_overwrite_requires_approval(settings):
    item = request("write_text", {"path": "output/result.txt", "content": "one"})
    execute(item, settings)
    with pytest.raises(ToolRejected, match="审批"):
        execute(item.model_copy(update={"arguments": {"path": "output/result.txt", "content": "two"}}), settings)
    approved = item.model_copy(update={"arguments": {"path": "output/result.txt", "content": "two"}, "approval_granted": True})
    assert execute(approved, settings).ok


def test_delete_moves_to_trash_and_requires_approval(settings):
    item = request("write_text", {"path": "workspace/a.txt", "content": "data"})
    execute(item, settings)
    deletion = item.model_copy(update={"tool": "move_to_trash", "arguments": {"path": "workspace/a.txt"}})
    with pytest.raises(ToolRejected, match="审批"):
        execute(deletion, settings)
    response = execute(deletion.model_copy(update={"approval_granted": True}), settings)
    assert response.data["path"].startswith("trash/")


def test_run_python_is_limited_to_python_file(settings):
    item = request("write_text", {"path": "workspace/job.py", "content": "print('ok')"})
    execute(item, settings)
    result = execute(item.model_copy(update={"tool": "run_python", "arguments": {"script": "workspace/job.py", "timeout_seconds": 5}}), settings)
    assert result.ok
    assert result.data["stdout"].strip() == "ok"


def test_run_python_normalizes_exact_exec_open_wrapper(settings):
    item = request("write_text", {"path": "workspace/job.py", "content": "print('normalized')"})
    execute(item, settings)
    result = execute(
        item.model_copy(
            update={
                "tool": "run_python",
                "arguments": {"script": "exec(open('workspace/job.py').read())"},
            }
        ),
        settings,
    )
    assert result.ok
    assert result.data["stdout"].strip() == "normalized"


def test_run_python_rejects_inline_code_with_actionable_message(settings):
    item = request("run_python", {"script": "print('not a path')"})
    with pytest.raises(ToolRejected, match="不能填写 Python 代码"):
        execute(item, settings)


def test_scoped_run_python_uses_task_root_relative_paths(settings):
    scope = {
        "node_key": "paper_2",
        "role": "document",
        "workspace": "workspace/agents/paper_2",
        "output": "output/paper_2",
        "readable_roots": [
            "input",
            "workspace/agents/paper_2",
            "shared/agents/paper_2",
            "output/paper_2",
        ],
        "writable_roots": [
            "workspace/agents/paper_2",
            "shared/agents/paper_2",
            "output/paper_2",
        ],
    }
    script = """\
import os
from pathlib import Path

Path("workspace/agents/paper_2/figs").mkdir(parents=True, exist_ok=True)
Path("workspace/agents/paper_2/figs/chart.txt").write_text("chart", encoding="utf-8")
Path("output/paper_2").mkdir(parents=True, exist_ok=True)
Path("output/paper_2/result.txt").write_text("ok", encoding="utf-8")
assert os.environ["MINISWARM_AGENT_WORKSPACE"] == "workspace/agents/paper_2"
assert os.environ["MINISWARM_AGENT_OUTPUT"] == "output/paper_2"
print("scoped-ok")
"""
    item = request(
        "write_text",
        {"path": "workspace/agents/paper_2/create_output.py", "content": script},
    ).model_copy(update={"agent_scope": scope})
    assert execute(item, settings).ok

    result = execute(
        item.model_copy(
            update={
                "tool": "run_python",
                "arguments": {
                    "script": "workspace/agents/paper_2/create_output.py",
                    "timeout_seconds": 10,
                },
            }
        ),
        settings,
    )
    assert result.ok, result.data["stderr"]
    assert result.data["stdout"].strip() == "scoped-ok"
    root = Path(settings.data_root) / "users" / str(item.user_id) / "tasks" / str(item.task_id)
    assert (root / "workspace/agents/paper_2/figs/chart.txt").read_text(encoding="utf-8") == "chart"
    assert (root / "output/paper_2/result.txt").read_text(encoding="utf-8") == "ok"


def test_document_libraries_create_and_verify_files(settings):
    script = r'''
from pathlib import Path
from docx import Document
from openpyxl import Workbook, load_workbook
from pptx import Presentation
from reportlab.pdfgen import canvas
from pypdf import PdfReader

out = Path("../output")
doc = Document(); doc.add_paragraph("MiniSwarm"); doc.save(out / "report.docx")
assert Document(out / "report.docx").paragraphs[0].text == "MiniSwarm"
wb = Workbook(); wb.active["A1"] = "MiniSwarm"; wb.save(out / "report.xlsx")
assert load_workbook(out / "report.xlsx").active["A1"].value == "MiniSwarm"
prs = Presentation(); prs.slides.add_slide(prs.slide_layouts[0]); prs.save(out / "report.pptx")
assert len(Presentation(out / "report.pptx").slides) == 1
pdf = canvas.Canvas(str(out / "report.pdf")); pdf.drawString(72, 720, "MiniSwarm"); pdf.save()
assert len(PdfReader(out / "report.pdf").pages) == 1
print("verified")
'''
    item = request("write_text", {"path": "workspace/documents.py", "content": script})
    execute(item, settings)
    result = execute(
        item.model_copy(update={"tool": "run_python", "arguments": {"script": "workspace/documents.py", "timeout_seconds": 30}}),
        settings,
    )
    assert result.ok, result.data["stderr"]
    assert result.data["stdout"].strip() == "verified"
    inspection = execute(
        item.model_copy(update={"tool": "inspect_document", "arguments": {"path": "output/report.pdf"}}),
        settings,
    )
    assert inspection.ok
    assert inspection.data["pages"] == 1


def test_inspect_document_rejects_empty_docx(settings):
    from docx import Document

    item = request("inspect_document", {"path": "output/empty.docx"})
    root = Path(settings.data_root) / "users" / str(item.user_id) / "tasks" / str(item.task_id)
    output = root / "output"
    output.mkdir(parents=True, exist_ok=True)
    Document().save(output / "empty.docx")
    with pytest.raises(ToolRejected, match="有效正文"):
        execute(item, settings)


def test_inspect_document_rejects_xlsx_formula_reference_error(settings):
    from openpyxl import Workbook

    item = request("inspect_document", {"path": "output/formula-error.xlsx"})
    root = Path(settings.data_root) / "users" / str(item.user_id) / "tasks" / str(item.task_id)
    output = root / "output"
    output.mkdir(parents=True, exist_ok=True)
    workbook = Workbook()
    workbook.active["A1"] = "=#REF!"
    workbook.save(output / "formula-error.xlsx")
    with pytest.raises(ToolRejected, match="公式错误"):
        execute(item, settings)


def test_inspect_pdf_rejects_visual_qa_issue(settings, monkeypatch):
    from reportlab.pdfgen import canvas

    item = request("inspect_document", {"path": "output/problem.pdf"})
    root = Path(settings.data_root) / "users" / str(item.user_id) / "tasks" / str(item.task_id)
    output = root / "output"
    output.mkdir(parents=True, exist_ok=True)
    pdf = canvas.Canvas(str(output / "problem.pdf"))
    pdf.drawString(72, 720, "content")
    pdf.save()
    monkeypatch.setattr(
        "runner_app.tools._pdf_visual_qa",
        lambda *args, **kwargs: {
            "available": True,
            "issues": [{"page": 1, "kind": "blank_page"}],
            "warnings": [],
        },
    )
    with pytest.raises(ToolRejected, match="视觉质检失败"):
        execute(item, settings)


def test_inspect_document_rejects_corrupt_file(settings):
    item = request("write_text", {"path": "output/broken.pdf", "content": "not a pdf"})
    execute(item, settings)
    with pytest.raises(ToolRejected, match="损坏"):
        execute(
            item.model_copy(update={"tool": "inspect_document", "arguments": {"path": "output/broken.pdf"}}),
            settings,
        )


def test_inspect_document_rejects_overlapping_pptx_text(settings):
    from pptx import Presentation
    from pptx.util import Inches

    item = request("inspect_document", {"path": "output/overlap.pptx"})
    root = Path(settings.data_root) / "users" / str(item.user_id) / "tasks" / str(item.task_id)
    output = root / "output"
    output.mkdir(parents=True, exist_ok=True)
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    first = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    first.text = "第一段正文"
    second = slide.shapes.add_textbox(Inches(1.5), Inches(1.2), Inches(5), Inches(1))
    second.text = "与第一段明显重叠"
    presentation.save(output / "overlap.pptx")

    with pytest.raises(ToolRejected, match="文本越界或重叠"):
        execute(item, settings)


def test_inspect_document_allows_minor_pptx_text_box_intersection(settings):
    from pptx import Presentation
    from pptx.util import Inches

    item = request("inspect_document", {"path": "output/minor-overlap.pptx"})
    root = Path(settings.data_root) / "users" / str(item.user_id) / "tasks" / str(item.task_id)
    output = root / "output"
    output.mkdir(parents=True, exist_ok=True)
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    first = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4.5), Inches(1))
    first.text = "标题框"
    second = slide.shapes.add_textbox(Inches(4.8), Inches(1), Inches(4), Inches(1))
    second.text = "相邻副标题框"
    presentation.save(output / "minor-overlap.pptx")

    result = execute(item, settings)
    assert result.ok
    assert result.data["layout_issues"] == []


def test_inspect_document_supports_html_csv_text_and_zip(settings):
    html = request(
        "write_text",
        {
            "path": "output/report.html",
            "content": "<!doctype html><html><head><title>Report</title></head><body><main>Result</main></body></html>",
        },
    )
    execute(html, settings)
    html_result = execute(
        html.model_copy(update={"tool": "inspect_document", "arguments": {"path": "output/report.html"}}),
        settings,
    )
    assert html_result.ok
    assert html_result.data["title"] == "Report"
    assert html_result.data["text_chars"] > 0

    csv_file = html.model_copy(
        update={"tool": "write_text", "arguments": {"path": "output/data.csv", "content": "name,value\na,1\n"}}
    )
    execute(csv_file, settings)
    csv_result = execute(
        csv_file.model_copy(update={"tool": "inspect_document", "arguments": {"path": "output/data.csv"}}),
        settings,
    )
    assert csv_result.data["rows_scanned"] == 2

    text_file = html.model_copy(
        update={"tool": "write_text", "arguments": {"path": "output/notes.md", "content": "# Notes\nVerified"}}
    )
    execute(text_file, settings)
    text_result = execute(
        text_file.model_copy(update={"tool": "inspect_document", "arguments": {"path": "output/notes.md"}}),
        settings,
    )
    assert text_result.data["nonempty_lines"] == 2

    archive = html.model_copy(
        update={
            "tool": "create_zip",
            "arguments": {"sources": ["output/notes.md"], "target": "output/result.zip"},
        }
    )
    execute(archive, settings)
    zip_result = execute(
        archive.model_copy(update={"tool": "inspect_document", "arguments": {"path": "output/result.zip"}}),
        settings,
    )
    assert zip_result.data["entries"] == 1


def test_run_tests_uses_fixed_pytest_command(settings):
    item = request(
        "write_text",
        {"path": "workspace/test_sample.py", "content": "def test_ok():\n    assert 2 + 2 == 4\n"},
    )
    execute(item, settings)
    result = execute(
        item.model_copy(
            update={
                "tool": "run_tests",
                "arguments": {"path": "workspace/test_sample.py", "timeout_seconds": 30},
            }
        ),
        settings,
    )
    assert result.ok, result.data["stderr"]
    assert "1 passed" in result.data["stdout"]
